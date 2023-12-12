from pptx import Presentation
from pptx.util import Inches, Pt
import fitz  # PyMuPDF

def pdf_to_pptx(pdf_path):
    # Open the PDF
    doc = fitz.open(pdf_path)

    # Create a PowerPoint presentation
    pptx = Presentation()

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap()
        image_path = f"/mnt/data/temp_page_{page_num}.png"
        pix.save(image_path)

        # Add a slide
        slide_layout = pptx.slide_layouts[5]  # choosing a blank slide layout
        slide = pptx.slides.add_slide(slide_layout)

        # Add the PDF page as an image to the slide
        left = Inches(1)
        top = Inches(1)
        slide.shapes.add_picture(image_path, left, top, width=pptx.slide_width - Inches(2))

    pptx_path = "/mnt/data/converted_presentation.pptx"
    pptx.save(pptx_path)
    return pptx_path

# Convert the PDF to PowerPoint
pdf_path = "/mnt/data/oreillytalkdec20231701450543622.pdf"
converted_pptx_path = pdf_to_pptx(pdf_path)
converted_pptx_path

